"""
Genera Eurisko-Sito-Presentazione-CAO.pptx per la presentazione interna al CAO.
Slide 16:9, palette navy/crimson coerente con il brand del sito.
Eseguire: python _build_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- Brand palette ---
NAVY = RGBColor(0x0A, 0x16, 0x28)
NAVY_SOFT = RGBColor(0x14, 0x24, 0x3D)
CRIMSON = RGBColor(0xC4, 0x1E, 0x3A)
CRIMSON_SOFT = RGBColor(0xFF, 0x4D, 0x66)
CREAM = RGBColor(0xF5, 0xF1, 0xE8)
INK_SOFT = RGBColor(0xC8, 0xC4, 0xBC)
LINE_FAINT = RGBColor(0x35, 0x40, 0x55)

# --- Setup presentation 16:9 ---
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
BLANK = prs.slide_layouts[6]  # blank layout


def add_bg(slide, color=NAVY):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_text(slide, text, left, top, width, height, *,
             font_size=18, bold=False, color=CREAM, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, italic=False, font_name="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    run.font.color.rgb = color
    return tb


def add_bullets(slide, bullets, left, top, width, height, *,
                font_size=20, color=CREAM, line_spacing=1.25):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(8)
        # Render "—" prefix in crimson, body in cream
        run_dash = p.add_run()
        run_dash.text = "— "
        run_dash.font.size = Pt(font_size)
        run_dash.font.color.rgb = CRIMSON_SOFT
        run_dash.font.bold = True
        run_dash.font.name = "Calibri"
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tb


def add_accent_bar(slide, left, top, width=Inches(0.5), height=Inches(0.06)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = CRIMSON
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def add_eyebrow(slide, text, left, top):
    return add_text(
        slide, text.upper(), left, top, Inches(6), Inches(0.4),
        font_size=12, color=CRIMSON_SOFT, bold=True
    )


def add_footer(slide, text="Eurisko S.r.l. · Sito istituzionale · Aprile 2026"):
    add_text(
        slide, text, Inches(0.6), Inches(7.05), Inches(12), Inches(0.3),
        font_size=10, color=INK_SOFT
    )


def add_slide_number(slide, n, total):
    add_text(
        slide, f"{n:02d} / {total:02d}", Inches(11.7), Inches(7.05),
        Inches(1.2), Inches(0.3),
        font_size=10, color=INK_SOFT, align=PP_ALIGN.RIGHT
    )


def add_speaker_notes(slide, notes):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = notes


# ===== build slides =====
slides_data = []

# 1. Cover
def slide_cover(s):
    add_bg(s, NAVY)
    # decorative crimson glow circle (ellipse) top-right
    glow = s.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(9.5), Inches(-1.5), Inches(5.5), Inches(5.5)
    )
    glow.fill.solid()
    glow.fill.fore_color.rgb = CRIMSON
    glow.line.fill.background()
    glow.shadow.inherit = False
    # bottom-left accent square
    sq = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(-1), Inches(5.5), Inches(3), Inches(3)
    )
    sq.fill.solid()
    sq.fill.fore_color.rgb = NAVY_SOFT
    sq.line.fill.background()
    sq.shadow.inherit = False

    add_eyebrow(s, "Eurisko · Presentazione interna", Inches(0.7), Inches(0.7))
    add_accent_bar(s, Inches(0.7), Inches(1.2), width=Inches(1.2), height=Inches(0.1))
    add_text(
        s, "Il nuovo sito istituzionale.",
        Inches(0.7), Inches(1.8), Inches(11), Inches(1.6),
        font_size=64, bold=True, color=CREAM
    )
    add_text(
        s, "Asset commerciale, vetrina di posizionamento, motore di lead generation\ne recruiting. Bilingue, accessibile, conforme.",
        Inches(0.7), Inches(3.6), Inches(11), Inches(1.5),
        font_size=22, color=INK_SOFT, italic=True
    )
    add_text(
        s, "Briefing per il CAO  ·  Aprile 2026",
        Inches(0.7), Inches(6.4), Inches(11), Inches(0.5),
        font_size=14, color=CRIMSON_SOFT, bold=True
    )

slides_data.append(("Cover", slide_cover, "Apri la presentazione introducendo il contesto: il nuovo sito istituzionale e gli obiettivi strategici. Tono diretto, concreto."))

# 2. Obiettivo e contesto
def slide_obiettivi(s):
    add_bg(s, NAVY)
    add_eyebrow(s, "01 · Obiettivo", Inches(0.7), Inches(0.6))
    add_accent_bar(s, Inches(0.7), Inches(1.0))
    add_text(
        s, "Perché abbiamo rifatto il sito.",
        Inches(0.7), Inches(1.3), Inches(12), Inches(1.2),
        font_size=44, bold=True, color=CREAM
    )
    add_bullets(s, [
        "Allineare la presenza online al posizionamento premium B2B di Eurisko nella consulenza SAP.",
        "Aprire il bacino oltre l'Italia con una versione inglese completa, non una traduzione di facciata.",
        "Trasformare il sito in uno strumento commerciale: lead-generation, recruiting, content marketing.",
        "Eliminare obsolescenza tecnica e dipendenze da CMS / plugin di terze parti.",
        "Rispetto pieno di GDPR, accessibilità WCAG 2.1 AA e best-practice SEO.",
    ], Inches(0.7), Inches(3.0), Inches(12), Inches(4), font_size=20)

slides_data.append(("Obiettivo", slide_obiettivi, "5 punti chiave. Insistere sul fatto che NON è solo una vetrina ma un asset commerciale. La versione EN apre mercato europeo per S/4HANA migration & AMS."))

# 3. I numeri
def slide_numeri(s):
    add_bg(s, NAVY)
    add_eyebrow(s, "02 · I numeri del progetto", Inches(0.7), Inches(0.6))
    add_accent_bar(s, Inches(0.7), Inches(1.0))
    add_text(
        s, "In sintesi.",
        Inches(0.7), Inches(1.3), Inches(12), Inches(1.2),
        font_size=44, bold=True, color=CREAM
    )

    stats = [
        ("72", "pagine pubblicate"),
        ("2", "lingue complete\nIT + EN"),
        ("4", "moduli SAP\nFI · CO · MM · SD"),
        ("10", "settori industriali\ncoperti"),
    ]
    col_w = Inches(2.8)
    gap = Inches(0.25)
    total_w = col_w * 4 + gap * 3
    start_x = (SLIDE_W - total_w) / 2
    for i, (num, label) in enumerate(stats):
        x = start_x + (col_w + gap) * i
        # card background
        card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(3.3), col_w, Inches(2.8))
        card.fill.solid()
        card.fill.fore_color.rgb = NAVY_SOFT
        card.line.color.rgb = LINE_FAINT
        card.line.width = Pt(0.75)
        card.shadow.inherit = False
        # number
        add_text(
            s, num, x, Inches(3.5), col_w, Inches(1.6),
            font_size=84, bold=True, color=CRIMSON_SOFT, align=PP_ALIGN.CENTER
        )
        # label
        add_text(
            s, label, x, Inches(5.1), col_w, Inches(1.0),
            font_size=14, color=INK_SOFT, align=PP_ALIGN.CENTER
        )

    add_text(
        s, "Tempo di sviluppo: 4 settimane · Costi infrastrutturali a regime: ~15 €/anno",
        Inches(0.7), Inches(6.5), Inches(12), Inches(0.4),
        font_size=14, color=CRIMSON_SOFT, italic=True, align=PP_ALIGN.CENTER
    )

slides_data.append(("Numeri", slide_numeri, "I numeri parlano da soli: 72 pagine = sito enterprise-grade, non vetrina. 4 moduli scelti = focus commerciale (FI/CO/MM/SD sono i piu' richiesti). 10 settori = leva SEO + credibilita' industriale."))

# 4. Cosa contiene il sito
def slide_struttura(s):
    add_bg(s, NAVY)
    add_eyebrow(s, "03 · Struttura del sito", Inches(0.7), Inches(0.6))
    add_accent_bar(s, Inches(0.7), Inches(1.0))
    add_text(
        s, "Cosa c'è dentro.",
        Inches(0.7), Inches(1.3), Inches(12), Inches(1.2),
        font_size=44, bold=True, color=CREAM
    )

    cols = [
        ("Cosa facciamo", [
            "Consulenza SAP",
            "Soluzioni per settore",
            "Migrazione S/4HANA",
            "AMS — Application Management",
        ]),
        ("Moduli SAP", [
            "Finance (FI)",
            "Controlling (CO)",
            "Material Management (MM)",
            "Sales & Distribution (SD)",
        ]),
        ("Settori industriali", [
            "Aerospazio · Automotive",
            "Chimica · Food & Beverage",
            "Industriale · Retail",
            "Public Sector · Travel",
            "Comm. & Media · Utilities",
        ]),
        ("Azienda", [
            "La nostra visione",
            "Portfolio progetti",
            "Lavora con noi",
            "Glossario SAP (80+ termini)",
            "Contatti",
        ]),
    ]
    col_w = Inches(2.95)
    gap = Inches(0.15)
    total_w = col_w * 4 + gap * 3
    start_x = (SLIDE_W - total_w) / 2
    for i, (title, items) in enumerate(cols):
        x = start_x + (col_w + gap) * i
        # title
        add_text(
            s, title, x, Inches(2.9), col_w, Inches(0.5),
            font_size=18, bold=True, color=CRIMSON_SOFT
        )
        # accent
        add_accent_bar(s, x, Inches(3.45), width=Inches(0.4), height=Inches(0.03))
        # items
        body = "\n".join(items)
        add_text(
            s, body, x, Inches(3.7), col_w, Inches(3.5),
            font_size=15, color=CREAM
        )

slides_data.append(("Struttura", slide_struttura, "Quattro colonne tematiche. Ogni voce è una pagina dedicata, non un placeholder. Glossario SAP = strategia content marketing per posizionamento organico."))

# 5. Focus: pagine settori e moduli
def slide_pagine(s):
    add_bg(s, NAVY)
    add_eyebrow(s, "04 · Profondità di contenuto", Inches(0.7), Inches(0.6))
    add_accent_bar(s, Inches(0.7), Inches(1.0))
    add_text(
        s, "Pagine verticali, non landing.",
        Inches(0.7), Inches(1.3), Inches(12), Inches(1.2),
        font_size=44, bold=True, color=CREAM
    )
    add_bullets(s, [
        "Ogni settore ha una pagina dedicata: dolori specifici, esempi concreti, moduli SAP rilevanti, casi d'uso.",
        "Ogni modulo SAP ha una pagina con panoramica funzionale, scenari tipici, deliverable Eurisko, CTA personalizzata.",
        "AMS e Migrazione S/4HANA come pagine-prodotto strategiche: sono i due servizi a maggior valore.",
        "Portfolio con case study anonimizzati ma concreti, divisi per settore.",
        "Glossario SAP con oltre 80 termini: strumento di posizionamento organico e supporto vendita.",
    ], Inches(0.7), Inches(3.0), Inches(12), Inches(4), font_size=20)

slides_data.append(("Profondita'", slide_pagine, "Sottolineare che ogni pagina e' contenutistica, non vuota. Pagina settore Automotive = ~800 parole rilevanti. Buon segnale per Google."))

# 6. Bilingue
def slide_bilingue(s):
    add_bg(s, NAVY)
    add_eyebrow(s, "05 · Strategia internazionale", Inches(0.7), Inches(0.6))
    add_accent_bar(s, Inches(0.7), Inches(1.0))
    add_text(
        s, "Bilingue completo, non superficiale.",
        Inches(0.7), Inches(1.3), Inches(12), Inches(1.2),
        font_size=44, bold=True, color=CREAM
    )
    add_bullets(s, [
        "Versione EN integralmente parallela a quella IT: 36 pagine ciascuna.",
        "Switch lingua su ogni pagina, mantiene il contesto (es. Cosa Facciamo ↔ What We Do).",
        "Hreflang dichiarato a Google: indicizzazione corretta delle coppie linguistiche.",
        "Traduzione adattata, non automatica: termini tecnici SAP allineati al lessico anglofono di settore.",
        "Apre il sito a clienti europei e gruppi multinazionali con HQ fuori Italia.",
    ], Inches(0.7), Inches(3.0), Inches(12), Inches(4), font_size=20)

slides_data.append(("Bilingue", slide_bilingue, "Investimento importante: la versione EN non e' un'opzione di facciata. Apre opportunita' di lead da multinazionali con sede europea fuori Italia."))

# 7. Lead generation: form
def slide_form(s):
    add_bg(s, NAVY)
    add_eyebrow(s, "06 · Lead generation & recruiting", Inches(0.7), Inches(0.6))
    add_accent_bar(s, Inches(0.7), Inches(1.0))
    add_text(
        s, "Due funnel attivi.",
        Inches(0.7), Inches(1.3), Inches(12), Inches(1.2),
        font_size=44, bold=True, color=CREAM
    )

    # two columns
    col_w = Inches(5.8)
    x1 = Inches(0.7)
    x2 = Inches(6.85)
    y = Inches(2.8)
    # col 1
    add_text(s, "Form contatti commerciali", x1, y, col_w, Inches(0.5),
             font_size=22, bold=True, color=CRIMSON_SOFT)
    add_accent_bar(s, x1, Inches(3.4), width=Inches(0.4), height=Inches(0.03))
    add_text(s,
             "• Campi essenziali · oggetto strutturato\n"
             "• Privacy / GDPR · consenso esplicito\n"
             "• Honeypot anti-bot Formspree\n"
             "• reCAPTCHA v2 Google\n"
             "• Pagina Grazie con redirect post-invio",
             x1, Inches(3.7), col_w, Inches(3.0),
             font_size=17, color=CREAM)
    # col 2
    add_text(s, "Form candidature", x2, y, col_w, Inches(0.5),
             font_size=22, bold=True, color=CRIMSON_SOFT)
    add_accent_bar(s, x2, Inches(3.4), width=Inches(0.4), height=Inches(0.03))
    add_text(s,
             "• Nome, Cognome, Città, Email, Telefono con prefisso\n"
             "• Posizione (collegata alle posizioni aperte)\n"
             "• Anni di esperienza SAP, LinkedIn\n"
             "• Permesso di lavoro · Categorie protette (L. 68/99)\n"
             "• Upload CV · GDPR · reCAPTCHA",
             x2, Inches(3.7), col_w, Inches(3.0),
             font_size=17, color=CREAM)

slides_data.append(("Form", slide_form, "Il sito non e' passivo: ha due canali attivi che alimentano commerciale e HR. Conformita' lavoristica italiana (categorie protette, permesso di lavoro) gestita nel form."))

# 8. Mobile
def slide_mobile(s):
    add_bg(s, NAVY)
    add_eyebrow(s, "07 · Esperienza utente", Inches(0.7), Inches(0.6))
    add_accent_bar(s, Inches(0.7), Inches(1.0))
    add_text(
        s, "Mobile-first, non mobile-after.",
        Inches(0.7), Inches(1.3), Inches(12), Inches(1.2),
        font_size=44, bold=True, color=CREAM
    )
    add_bullets(s, [
        "Layout responsive su 3 breakpoint: desktop, tablet, mobile.",
        "Menu hamburger con accordion per i sotto-livelli.",
        "Animazioni cinematiche su hero (rivelazione progressiva, glow, parallasse).",
        "Disabilitazione automatica delle animazioni se l'utente le ha disattivate dal sistema operativo (accessibilità).",
        "Performance: CSS e JS minificati, immagini ottimizzate, lazy-loading.",
        "Search interna con scorciatoie da tastiera (ESC per chiudere, Ctrl+Click per nuova scheda).",
    ], Inches(0.7), Inches(3.0), Inches(12), Inches(4), font_size=18)

slides_data.append(("Mobile/UX", slide_mobile, "Il 60-70% del traffico B2B passa da mobile (smartphone in pausa caffe', tablet in riunione). Performance reale > glamour della homepage."))

# 9. Architettura tecnica
def slide_tecnico(s):
    add_bg(s, NAVY)
    add_eyebrow(s, "08 · Architettura tecnica", Inches(0.7), Inches(0.6))
    add_accent_bar(s, Inches(0.7), Inches(1.0))
    add_text(
        s, "Stack semplice. Affidabile. Senza lock-in.",
        Inches(0.7), Inches(1.3), Inches(12), Inches(1.2),
        font_size=40, bold=True, color=CREAM
    )

    rows = [
        ("Sorgente", "HTML + CSS + JavaScript vanilla. No framework, no CMS, no database."),
        ("Hosting", "Vercel · CDN globale · deploy automatico da Git · HTTPS gratuito"),
        ("Versionamento", "Repository GitHub privato · cronologia completa · backup gratuito"),
        ("Form & e-mail", "Formspree gestisce le submissions e inoltra alle mailbox interne"),
        ("Anti-spam", "reCAPTCHA Google + honeypot Formspree + validazione client"),
        ("Analytics", "Predisposto per Google Analytics / Search Console / privacy-friendly"),
    ]
    y = Inches(2.95)
    for label, val in rows:
        add_text(s, label, Inches(0.7), y, Inches(2.6), Inches(0.4),
                 font_size=15, bold=True, color=CRIMSON_SOFT)
        add_text(s, val, Inches(3.4), y, Inches(9.5), Inches(0.4),
                 font_size=15, color=CREAM)
        y += Inches(0.6)

slides_data.append(("Tecnico", slide_tecnico, "Insistere sul concetto NO LOCK-IN. Se Vercel sparisse, il sito si sposta su qualsiasi hosting in mezz'ora. I form si ricollegano a un altro servizio. Costi e dipendenze minimi."))

# 10. SEO & Marketing
def slide_seo(s):
    add_bg(s, NAVY)
    add_eyebrow(s, "09 · SEO & content marketing", Inches(0.7), Inches(0.6))
    add_accent_bar(s, Inches(0.7), Inches(1.0))
    add_text(
        s, "Indicizzazione costruita on-page.",
        Inches(0.7), Inches(1.3), Inches(12), Inches(1.2),
        font_size=42, bold=True, color=CREAM
    )
    add_bullets(s, [
        "Meta description, Open Graph e Twitter Card configurati su ogni pagina.",
        "Sitemap XML + robots.txt + structured data JSON-LD (breadcrumb).",
        "Hreflang IT/EN per indicare a Google le coppie di pagine equivalenti.",
        "Keyword strategy: \"consulenza SAP\", \"migrazione S/4HANA\", \"AMS SAP\", per modulo, per settore.",
        "Glossario SAP (80+ termini) come asset organico per long-tail.",
        "Performance Lighthouse target 90+ su SEO, accessibilità e best practice.",
    ], Inches(0.7), Inches(3.0), Inches(12), Inches(4), font_size=18)

slides_data.append(("SEO", slide_seo, "Risultati SEO si misurano a 30/60/90 giorni dal go-live. Il sito e' tecnicamente pronto, manca solo il via in Search Console + dominio puntato."))

# 11. Accessibilita
def slide_a11y(s):
    add_bg(s, NAVY)
    add_eyebrow(s, "10 · Accessibilità & GDPR", Inches(0.7), Inches(0.6))
    add_accent_bar(s, Inches(0.7), Inches(1.0))
    add_text(
        s, "Conformità non opzionale.",
        Inches(0.7), Inches(1.3), Inches(12), Inches(1.2),
        font_size=44, bold=True, color=CREAM
    )

    col_w = Inches(5.8)
    x1 = Inches(0.7)
    x2 = Inches(6.85)
    y = Inches(2.8)
    # Accessibilita
    add_text(s, "Accessibilità — WCAG 2.1 AA", x1, y, col_w, Inches(0.5),
             font_size=20, bold=True, color=CRIMSON_SOFT)
    add_accent_bar(s, x1, Inches(3.4), width=Inches(0.4), height=Inches(0.03))
    add_text(s,
             "• Skip link, navigazione completa da tastiera\n"
             "• Aria-label e ruoli semantici\n"
             "• Contrasti conformi · alt text su tutte le immagini\n"
             "• Rispetto preferenze sistema (riduzione movimento)\n"
             "• Pagina Dichiarazione di Accessibilità",
             x1, Inches(3.7), col_w, Inches(3.0),
             font_size=16, color=CREAM)
    # GDPR
    add_text(s, "GDPR / Compliance", x2, y, col_w, Inches(0.5),
             font_size=20, bold=True, color=CRIMSON_SOFT)
    add_accent_bar(s, x2, Inches(3.4), width=Inches(0.4), height=Inches(0.03))
    add_text(s,
             "• Cookie banner conforme · consenso opt-in granulare\n"
             "• Privacy Policy · Cookie Policy · Note Legali\n"
             "• Termini di Utilizzo · Mappa del Sito\n"
             "• Consenso esplicito su form (commerciale + selezione)\n"
             "• Nessun tracking di terze parti senza consenso",
             x2, Inches(3.7), col_w, Inches(3.0),
             font_size=16, color=CREAM)

slides_data.append(("Compliance", slide_a11y, "Accessibilita' WCAG 2.1 AA = obbligo di legge per soggetti privati con servizio pubblico (Direttiva 2102 + Stanca). GDPR = base legale per ogni form. Pronti per audit."))

# 12. Costi
def slide_costi(s):
    add_bg(s, NAVY)
    add_eyebrow(s, "11 · Costi & manutenzione", Inches(0.7), Inches(0.6))
    add_accent_bar(s, Inches(0.7), Inches(1.0))
    add_text(
        s, "Spesa annua: irrisoria.",
        Inches(0.7), Inches(1.3), Inches(12), Inches(1.2),
        font_size=44, bold=True, color=CREAM
    )

    rows = [
        ("Dominio (1 anno)", "~ 15 €", "Rinnovo annuale"),
        ("Hosting Vercel", "0 €", "Piano free sufficiente per traffico B2B"),
        ("GitHub", "0 €", "Repository privato gratuito"),
        ("Formspree", "0 € · oppure ~10 €/mese", "Free fino a 50 submissions/mese"),
        ("Google reCAPTCHA", "0 €", "Servizio gratuito"),
        ("Totale a regime", "~ 15 € / anno", "+ Formspree solo se i volumi salgono"),
    ]
    y = Inches(2.85)
    for i, (label, price, note) in enumerate(rows):
        is_last = (i == len(rows) - 1)
        # row background for total
        if is_last:
            bg = s.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.55), y - Inches(0.05),
                Inches(12.2), Inches(0.6)
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = NAVY_SOFT
            bg.line.fill.background()
            bg.shadow.inherit = False
        color_label = CRIMSON_SOFT if is_last else CREAM
        weight = True if is_last else False
        add_text(s, label, Inches(0.7), y, Inches(4.5), Inches(0.4),
                 font_size=16, bold=weight, color=color_label)
        add_text(s, price, Inches(5.4), y, Inches(3.0), Inches(0.4),
                 font_size=16, bold=True, color=CRIMSON_SOFT)
        add_text(s, note, Inches(8.6), y, Inches(4.3), Inches(0.4),
                 font_size=14, color=INK_SOFT, italic=True)
        y += Inches(0.6)

slides_data.append(("Costi", slide_costi, "Numero che impressiona ogni CFO/CAO: ~15 euro/anno fissi. Manutenzione gestita internamente via repository Git. Nessun canone, nessuna dipendenza."))

# 13. Roadmap
def slide_roadmap(s):
    add_bg(s, NAVY)
    add_eyebrow(s, "12 · Prossimi passi", Inches(0.7), Inches(0.6))
    add_accent_bar(s, Inches(0.7), Inches(1.0))
    add_text(
        s, "Cosa succede dopo il go-live.",
        Inches(0.7), Inches(1.3), Inches(12), Inches(1.2),
        font_size=42, bold=True, color=CREAM
    )

    phases = [
        ("Subito", "0-30 gg", [
            "Puntamento dominio definitivo",
            "Indicizzazione su Google Search Console e Bing Webmaster",
            "Test finali su browser e device reali",
        ]),
        ("Breve termine", "1-3 mesi", [
            "Pubblicazione 2-3 case study reali (con consenso clienti)",
            "Primo articolo di approfondimento (es. \"5 errori in S/4HANA migration\")",
            "Monitoraggio analytics e prime ottimizzazioni",
        ]),
        ("Medio-lungo", "3-12 mesi", [
            "A/B test su CTA homepage, sezione headline",
            "Espansione glossario e sezione insight",
            "Eventuale newsletter / integrazione LinkedIn jobs",
        ]),
    ]
    col_w = Inches(4.05)
    gap = Inches(0.15)
    total_w = col_w * 3 + gap * 2
    start_x = (SLIDE_W - total_w) / 2
    for i, (title, when, items) in enumerate(phases):
        x = start_x + (col_w + gap) * i
        # card
        card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.85),
                                   col_w, Inches(3.9))
        card.fill.solid()
        card.fill.fore_color.rgb = NAVY_SOFT
        card.line.color.rgb = LINE_FAINT
        card.line.width = Pt(0.75)
        card.shadow.inherit = False
        # phase title
        add_text(s, title, x + Inches(0.25), Inches(3.0), col_w - Inches(0.5),
                 Inches(0.5), font_size=20, bold=True, color=CRIMSON_SOFT)
        add_text(s, when, x + Inches(0.25), Inches(3.5), col_w - Inches(0.5),
                 Inches(0.4), font_size=12, color=INK_SOFT, italic=True)
        # items
        body = "\n".join("• " + it for it in items)
        add_text(s, body, x + Inches(0.25), Inches(4.05), col_w - Inches(0.5),
                 Inches(2.8), font_size=14, color=CREAM)

slides_data.append(("Roadmap", slide_roadmap, "Sito gia' pronto. Prossimo step e' un'azione di 5 minuti (cambio DNS dominio). Tutto il resto e' incrementale e dipende dalle priorita' commerciali."))

# 14. Domande probabili
def slide_qa(s):
    add_bg(s, NAVY)
    add_eyebrow(s, "13 · Q&A — domande probabili", Inches(0.7), Inches(0.6))
    add_accent_bar(s, Inches(0.7), Inches(1.0))
    add_text(
        s, "Possibili obiezioni e risposte.",
        Inches(0.7), Inches(1.3), Inches(12), Inches(1.2),
        font_size=40, bold=True, color=CREAM
    )

    qa = [
        ("Quanto siamo competitivi sul SEO?",
         "Da misurare a 30/60/90 giorni. Content è già denso e localizzato per le keyword target."),
        ("Possiamo modificare i testi senza chiamare uno sviluppatore?",
         "Sì, modifica diretta del file + push su GitHub. Posso predisporre una guida operativa."),
        ("Cosa succede se Vercel o Formspree chiudono?",
         "Sito statico → trasferibile su qualsiasi hosting in mezz'ora. Form ricollegabili a servizi alternativi (Resend, Mailgun, formcarry). Zero lock-in."),
        ("I dati dei candidati sono al sicuro?",
         "Trasferimento HTTPS, nessun database lato nostro: i CV arrivano alla mailbox HR. Consenso GDPR esplicito, finalità delimitata, conservazione regolata da informativa."),
        ("Si può aggiungere un'area riservata clienti?",
         "Sì, ma è un progetto a parte (autenticazione + backend). Da valutare in base al beneficio."),
    ]
    y = Inches(2.85)
    for q, a in qa:
        add_text(s, "?  " + q, Inches(0.7), y, Inches(12), Inches(0.4),
                 font_size=15, bold=True, color=CRIMSON_SOFT)
        add_text(s, a, Inches(0.95), y + Inches(0.4), Inches(11.8), Inches(0.5),
                 font_size=14, color=CREAM, italic=True)
        y += Inches(0.85)

slides_data.append(("Q&A", slide_qa, "Cinque domande tipiche con risposte gia' pronte. Se il CAO ne fa altre, prendere appunti e rispondere a posteriori."))

# 15. Closing
def slide_closing(s):
    add_bg(s, NAVY)
    # decorative crimson glow circle (ellipse) bottom-left
    glow = s.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(-2.5), Inches(4), Inches(6), Inches(6)
    )
    glow.fill.solid()
    glow.fill.fore_color.rgb = CRIMSON
    glow.line.fill.background()
    glow.shadow.inherit = False
    # accent square top right
    sq = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(11), Inches(-1), Inches(3), Inches(3)
    )
    sq.fill.solid()
    sq.fill.fore_color.rgb = NAVY_SOFT
    sq.line.fill.background()
    sq.shadow.inherit = False

    add_eyebrow(s, "Eurisko", Inches(0.7), Inches(2.0))
    add_accent_bar(s, Inches(0.7), Inches(2.5), width=Inches(1.2), height=Inches(0.1))
    add_text(
        s, "Smart, functional, dynamic,\nproactive, agile.",
        Inches(0.7), Inches(2.9), Inches(12), Inches(2.2),
        font_size=56, bold=True, color=CREAM
    )
    add_text(
        s, "Il tuo partner SAP.",
        Inches(0.7), Inches(5.3), Inches(12), Inches(0.7),
        font_size=22, color=CRIMSON_SOFT, italic=True
    )
    add_text(
        s, "Grazie.  ·  Domande?",
        Inches(0.7), Inches(6.5), Inches(12), Inches(0.5),
        font_size=14, color=INK_SOFT
    )

slides_data.append(("Chiusura", slide_closing, "Slide finale. Tagline aziendale come chiusura emotiva. Lasciare la slide visibile durante il Q&A."))


# Build all slides
total = len(slides_data)
for i, (title, builder, notes) in enumerate(slides_data, start=1):
    slide = prs.slides.add_slide(BLANK)
    builder(slide)
    if i not in (1, total):  # cover and closing: no number/footer
        add_slide_number(slide, i, total)
        add_footer(slide)
    add_speaker_notes(slide, notes)


out = "Eurisko-Sito-Presentazione-CAO.pptx"
prs.save(out)
print(f"Generato: {out} · {total} slide")
