"""Genera un file pptx con la sola slide 'Identita' visiva' da inserire poi
nel deck originale Eurisko-Sito-Presentazione-CAO.pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY = RGBColor(0x0A, 0x16, 0x28)
NAVY_SOFT = RGBColor(0x14, 0x24, 0x3D)
CRIMSON = RGBColor(0xC4, 0x1E, 0x3A)
CRIMSON_SOFT = RGBColor(0xFF, 0x4D, 0x66)
CREAM = RGBColor(0xF5, 0xF1, 0xE8)
INK_SOFT = RGBColor(0xC8, 0xC4, 0xBC)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W, SLIDE_H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_bg(slide, color=NAVY):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False


def add_text(slide, text, left, top, width, height, *,
             font_size=18, bold=False, color=CREAM, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, italic=False, font_name="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    run.font.color.rgb = color


def add_accent_bar(slide, left, top, width=Inches(0.5), height=Inches(0.06)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = CRIMSON
    bar.line.fill.background()
    bar.shadow.inherit = False


def add_eyebrow(slide, text, left, top):
    add_text(slide, text.upper(), left, top, Inches(8), Inches(0.4),
             font_size=12, color=CRIMSON_SOFT, bold=True)


# ===== Slide unica =====
slide = prs.slides.add_slide(BLANK)
add_bg(slide, NAVY)

add_eyebrow(slide, "09 · Identità visiva", Inches(0.7), Inches(0.6))
add_accent_bar(slide, Inches(0.7), Inches(1.0))

add_text(slide, "Premium nei dettagli.\nRiconoscibile a colpo d'occhio.",
         Inches(0.7), Inches(1.3), Inches(12), Inches(2.0),
         font_size=40, bold=True, color=CREAM)

add_text(slide,
         "Niente template, niente accent generici. Il sito è progettato come un editorial design.",
         Inches(0.7), Inches(3.05), Inches(12), Inches(0.5),
         font_size=16, color=INK_SOFT, italic=True)

# Bullets: lead bold + body
bullets = [
    ("Palette dedicata.",
     "Navy profondo come fondo, crimson come unico accent dosato con parsimonia. Tono autorevole, distintivo, immediatamente riconoscibile."),
    ("Tipografia editoriale.",
     "Switzer + Oswald + Space Grotesk: gerarchia da rivista di settore, non font di sistema da landing page."),
    ("Animazioni cinematiche.",
     "Reveal progressivi, tilt 3D parallasse, glow burst, particle footer interattivo. Wow factor funzionale, disattivabile dal sistema."),
    ("Iconografia custom.",
     "32 SVG line-art disegnati su misura per moduli SAP, settori e pagine. Stile coerente, peso del tratto stratificato. Nessun set stock."),
    ("Layout editoriale asimmetrico.",
     "Griglie a proporzioni non banali, divisori sottili, padding responsivi via CSS subgrid. Non un template Wix o Webflow."),
    ("Coerenza su 72 pagine.",
     "Stesso linguaggio visivo dalla home all'ultima pagina del glossario. Brand riconoscibile da qualsiasi punto di atterraggio."),
]

y = Inches(3.75)
row_h = Inches(0.52)
for lead, body in bullets:
    # crimson dash + bold lead
    tb1 = slide.shapes.add_textbox(Inches(0.7), y, Inches(4.6), row_h)
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_right = tf1.margin_top = tf1.margin_bottom = 0
    p1 = tf1.paragraphs[0]
    r_dash = p1.add_run()
    r_dash.text = "— "
    r_dash.font.size = Pt(15)
    r_dash.font.bold = True
    r_dash.font.color.rgb = CRIMSON_SOFT
    r_dash.font.name = "Calibri"
    r_lead = p1.add_run()
    r_lead.text = lead
    r_lead.font.size = Pt(15)
    r_lead.font.bold = True
    r_lead.font.color.rgb = CREAM
    r_lead.font.name = "Calibri"
    # body
    add_text(slide, body, Inches(5.4), y, Inches(7.5), row_h,
             font_size=14, color=INK_SOFT)
    y += row_h

# Closing footer line
add_text(slide,
         "Risultato percepito: trust enterprise. Per chi valuta a chi affidare una migrazione S/4HANA, il sito comunica già un livello di cura del dettaglio.",
         Inches(0.7), Inches(6.95), Inches(12), Inches(0.4),
         font_size=11, color=CRIMSON_SOFT, italic=True)

# Speaker notes
slide.notes_slide.notes_text_frame.text = (
    "Slide piu' soft del deck ma quella che lascia il segno emotivo. "
    "Un sito ben fatto suggerisce processi ben fatti. La maggior parte dei "
    "competitor SAP italiani ha siti datati o basati su template: chi atterra "
    "su euriskosrl.it percepisce una differenza di livello. Se il CAO chiede "
    "\"quanto fa risparmiare un template?\", risposta: \"il template costa "
    "meno in giornate ma costa di piu' in posizionamento perso\"."
)

out = "Eurisko-Slide-Identita-Visiva.pptx"
prs.save(out)
print(f"Generato: {out}")
